"""md2ppt — reusable pptx composition helpers.

Used by build scripts produced via the md2ppt skill.
Style preset: corporate_blue (default; depends on PingFang TC font installed).

Import in your build script:
    import sys
    sys.path.insert(0, "/Users/<you>/.claude/skills/md2ppt/scripts")
    from pptx_helpers import *

Then proceed with `prs = init_deck()` and `add_*` calls.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# Style preset: corporate_blue
# ============================================================
FONT       = "PingFang TC"
FONT_MONO  = "Menlo"
COLOR_BG     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TITLE  = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_TEXT   = RGBColor(0x21, 0x21, 0x21)
COLOR_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
COLOR_OK     = RGBColor(0x27, 0xAE, 0x60)
COLOR_WARN   = RGBColor(0xE6, 0x7E, 0x22)
COLOR_MUTED  = RGBColor(0x7F, 0x8C, 0x8D)
COLOR_BAR    = RGBColor(0x34, 0x98, 0xDB)


def init_deck(width_in=13.333, height_in=7.5):
    """Create a 16:9 Presentation. Default 13.333 x 7.5 inches."""
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs


def init_deck_from_template(template_path):
    """Open a brand template, strip its existing slides, return the Presentation.

    Inherits theme + slide_master + slide_layouts. Strips all existing slides
    so the build script starts from a clean deck but with brand chrome.

    Args:
        template_path: path to brand template .pptx

    Returns:
        Presentation object ready to add fresh slides

    Raises:
        FileNotFoundError if template_path doesn't exist
    """
    from pathlib import Path
    p = Path(template_path).expanduser()
    if not p.exists():
        raise FileNotFoundError(f"Brand template not found: {p}")
    prs = Presentation(str(p))

    # Strip all existing slides — keep only theme/masters/layouts.
    # python-pptx doesn't expose a clean .delete_slide(), so we manipulate
    # the underlying XML directly + use Part.drop_rel(). See:
    # https://github.com/scanny/python-pptx/issues/67
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    for sldId in sldIds:
        rId = sldId.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)
    return prs


def _find_layout(prs, layout_name_hint):
    """Find a slide layout by name (substring match, case-insensitive).

    Args:
        prs: Presentation
        layout_name_hint: substring to match against layout names

    Returns:
        SlideLayout matching the hint, or layout[6] (typically blank) as fallback
    """
    hint = layout_name_hint.lower()
    for layout in prs.slide_layouts:
        if hint in layout.name.lower():
            return layout
    # fallback: try blank-like
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or "空白" in layout.name:
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def add_slide_with_template_layout(prs, layout_name_hint):
    """Add a slide using the named template layout (or fallback). Returns the slide.

    Used internally by add_cover_from_template etc. Caller fills placeholders
    via slide.placeholders[N].text = '...' or by adding additional shapes.
    """
    layout = _find_layout(prs, layout_name_hint)
    slide = prs.slides.add_slide(layout)
    slide._prs = prs  # convenience back-ref for downstream helpers
    return slide


def add_cover_from_template(prs, *, layout_name="標題投影片", title=None, subtitle=None):
    """Cover slide using template's title-slide layout.

    Fills placeholder[0] with title and placeholder[1] (if present) with subtitle.
    Falls back to add_blank_slide + manual textbox if layout missing.

    Args:
        prs: Presentation
        layout_name: substring match for template layout (default "標題投影片")
        title: cover title (str, required for non-empty cover)
        subtitle: cover subtitle (str, optional)

    Returns:
        Slide
    """
    slide = add_slide_with_template_layout(prs, layout_name)
    placeholders = list(slide.placeholders)
    if title and len(placeholders) > 0:
        placeholders[0].text = title
    if subtitle and len(placeholders) > 1:
        placeholders[1].text = subtitle
    return slide


def add_blank_from_template(prs, *, layout_name="空白", clear_placeholders=True):
    """Add a slide using a template layout for content composition.

    Despite the name, this isn't restricted to "blank" layouts — pass any
    layout_name. Common patterns:
    - layout_name="空白" → minimal layout (no chrome, no placeholders)
    - layout_name="標題投影片" → layout with brand chrome (logo / page number /
      background images) inherited from layout shapes; placeholders cleared
      by default so they don't show "click to add title" prompts
    - layout_name="標題及內容" → similar but with explicit title + content
      placeholder (constraining)

    Use in brand mode for content slides where you want template's theme
    inheritance but full hand-coded layout freedom.

    **CRITICAL**: Does NOT add a white background rectangle. The template's
    layout/master typically has its own background graphics, logo, page number,
    or footer chrome. Adding a white rect would hide them all. If you need
    a clean white section for content, place white textbox/rect inside the
    content area only — never full-slide cover.

    Args:
        prs: Presentation (loaded via init_deck_from_template)
        layout_name: substring match for layout (default "空白")
        clear_placeholders: if True, set all inherited placeholders to empty
            string so they don't show default "click to add..." text. Default True.

    Returns:
        Slide
    """
    layout = _find_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)
    if clear_placeholders:
        for ph in slide.placeholders:
            try:
                ph.text = ""
            except Exception:
                pass
    slide._prs = prs
    return slide


def add_section_divider_from_template(prs, *, layout_name="章節標題", title=None):
    """Section divider slide using template's section-header layout.

    Args:
        prs: Presentation
        layout_name: substring match for template layout (default "章節標題")
        title: section title

    Returns:
        Slide
    """
    slide = add_slide_with_template_layout(prs, layout_name)
    placeholders = list(slide.placeholders)
    if title and len(placeholders) > 0:
        placeholders[0].text = title
    return slide


def list_template_layouts(template_path):
    """Inspect a template — print available layouts. Useful at design time.

    Args:
        template_path: path to .pptx
    """
    prs = Presentation(str(template_path))
    print(f"Template: {template_path}")
    print(f"Layouts ({len(prs.slide_layouts)}):")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")


def add_blank_slide(prs):
    """Add a blank slide with white background. Returns the slide."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    slide._prs = prs  # convenience back-ref so helpers can read slide_width/height
    return slide


def add_textbox(slide, left, top, width, height, text, *,
                font=FONT, size=18, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, line_spacing=1.15):
    """Plain text frame. `text` can be string (with \\n) or list of lines."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_title_bar(slide, title, subtitle=None):
    """Standard top title bar: 30pt bold title + optional 14pt muted subtitle + thin underline."""
    prs = slide._prs
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7),
                title, size=30, bold=True, color=COLOR_TITLE)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.4),
                    subtitle, size=14, color=COLOR_MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(1.45),
                                   Inches(12.1), Emu(20000))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TITLE
    line.line.fill.background()


def add_bullets(slide, left, top, width, height, items, *, size=16, bold=False, color=COLOR_TEXT):
    """Bullet list. items = list of (text, level) or just text strings (level=0)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, lvl = item
        else:
            text, lvl = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.level = lvl
        bullet = "• " if lvl == 0 else "  – "
        run = p.add_run()
        run.text = bullet + text
        run.font.name = FONT
        run.font.size = Pt(size - lvl * 2)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_mono_block(slide, left, top, width, height, text, *,
                   size=12, color=COLOR_TEXT, highlight_rules=None):
    """Monospace text frame for ASCII art / code. `highlight_rules`: list of (substr, color) tuples."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = line
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        c = color
        if highlight_rules:
            for rule_match, rule_color in highlight_rules:
                if rule_match in line:
                    c = rule_color
                    break
        run.font.color.rgb = c
    return tb


def add_table(slide, left, top, width, height, data, *,
              header_color=COLOR_TITLE, font_size=14, col_widths=None):
    """Table with zebra striping. data = list of rows (each row = list of cells).

    `col_widths`: optional list of Inches(N) — sets explicit per-column width.
    """
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = FONT
            run.font.size = Pt(font_size)
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                run.font.color.rgb = COLOR_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 1 else RGBColor(0xF4, 0xF6, 0xF8)
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    return tbl


def add_picture_fit(slide, png_path, *, top, max_width, max_height,
                    vertical_center_in=None):
    """Add a picture with auto-fit (preserve aspect ratio) + horizontal center.

    `vertical_center_in`: tuple (top_bound, bottom_bound). If given, picture is
    vertically centered within that band (overrides `top`).
    """
    from PIL import Image
    prs = slide._prs
    iw, ih = Image.open(str(png_path)).size
    aspect = iw / ih
    h = max_height
    w = int(h * aspect)
    if w > max_width:
        w = max_width
        h = int(w / aspect)
    left = int((prs.slide_width - w) / 2)
    if vertical_center_in is not None:
        tb, bb = vertical_center_in
        top = tb + (bb - tb - h) // 2
    return slide.shapes.add_picture(str(png_path), left, top, width=w, height=h)


def add_log_bar(slide, label, value_kbps, value_label, *,
                top, color=COLOR_BAR, label_size=12,
                bar_left=None, bar_max_width=None, bar_height=None,
                max_kbps=800000):
    """Log-scaled horizontal bar for throughput/quantity comparisons.

    Positions:
    - left label: right-aligned in left column
    - bar:        starts at bar_left
    - value label: just right of bar end

    Defaults assume 13.333" wide slide.
    """
    import math
    prs = slide._prs
    if bar_left is None:
        bar_left = Inches(3.8)
    if bar_max_width is None:
        bar_max_width = Inches(7.5)
    if bar_height is None:
        bar_height = Inches(0.55)

    def log_w(v):
        if v <= 0:
            return 0.05
        return min(1.0, math.log10(max(1, v)) / math.log10(max_kbps))

    width = int(bar_max_width * log_w(value_kbps))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, top, width, bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, Inches(0.4), top, bar_left - Inches(0.5), bar_height,
                label, size=label_size, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)
    add_textbox(slide, bar_left + width + Inches(0.1), top, Inches(2.3), bar_height,
                value_label, size=label_size, color=COLOR_TEXT, bold=True)


def save(prs, path):
    """Save deck and print confirmation."""
    prs.save(path)
    print(f"OK → {path}")
    print(f"Slides: {len(prs.slides)}")
